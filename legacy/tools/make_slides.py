"""Build 6-slide deck for PhysMorph-GS SCA supp video."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGS = Path("/home/chayo/Desktop/Shape-morphing-binder/figs/paper_final")
OUT  = Path("/home/chayo/Desktop/Shape-morphing-binder/assets/PhysMorphGS_slides.pptx")

# Palette
INK     = RGBColor(0x1F, 0x2D, 0x3D)
MUTE    = RGBColor(0x5A, 0x6B, 0x7E)
ACCENT  = RGBColor(0xC0, 0x39, 0x2B)   # red
ACCENT2 = RGBColor(0x2E, 0x6A, 0xA6)   # blue
BG      = RGBColor(0xFC, 0xFC, 0xFA)
RULE    = RGBColor(0xDC, 0xDC, 0xD4)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]


def add_bg(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_header(slide, kicker, title):
    # Top rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(0.55), Inches(0.35), Inches(0.05))
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(1.05), Inches(0.40), Inches(8), Inches(0.35),
             kicker, size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.78), Inches(12.1), Inches(0.8),
             title, size=26, bold=True, color=INK)
    # Footer
    add_text(slide, Inches(0.6), Inches(7.1), Inches(8), Inches(0.3),
             "PhysMorph-GS  \u00b7  Render-Guided Volumetric Morphing with Differentiable Physics",
             size=10, color=MUTE)


def add_bullets(slide, x, y, w, h, items, *, size=16, gap=0.35, color=INK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        bold = False
        text = item
        if isinstance(item, tuple):
            text, bold = item
        r = p.add_run()
        r.text = "\u2022  " + text
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_image_fit(slide, path, x, y, w, h):
    # Keep aspect ratio centered in box
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = (w / h) if isinstance(w, (int, float)) else (w.emu / h.emu)
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        fit_w = w
        fit_h = Emu(int(w.emu / img_ratio))
    else:
        fit_h = h
        fit_w = Emu(int(h.emu * img_ratio))
    fx = x + Emu(int((w.emu - fit_w.emu) / 2))
    fy = y + Emu(int((h.emu - fit_h.emu) / 2))
    slide.shapes.add_picture(str(path), fx, fy, fit_w, fit_h)


# =====================================================================
# Slide 1 — Title + Problem
# =====================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "SCA 2026  \u00b7  Paper 1105", "Problem: render-guided volumetric morphing")

add_text(s, Inches(0.6), Inches(1.65), Inches(6.6), Inches(0.4),
         "Goal", size=15, bold=True, color=ACCENT2)
add_text(s, Inches(0.6), Inches(2.05), Inches(6.6), Inches(1.2),
         "Produce a physically plausible animation that deforms a source volume into a target shape \u2014 respecting material response while matching the target\u2019s visual silhouette.",
         size=15, color=INK)

add_text(s, Inches(0.6), Inches(3.25), Inches(6.6), Inches(0.4),
         "Why it\u2019s hard", size=15, bold=True, color=ACCENT2)
add_bullets(s, Inches(0.6), Inches(3.7), Inches(6.6), Inches(3.0), [
    ("Physics-only mass matching \u2192 coarse shape only.", True),
    "Grid-based loss cannot resolve thin protrusions, concavities, fine features.",
    ("Naive render coupling on positions \u2192 unstable.", True),
    "\u2202L/\u2202x fights elastic restoring forces and needs ad-hoc cross-space scaling.",
    "Eigen-decomposition failures at gains that F-space absorbs without issue.",
], size=14)

# Teaser image on right
add_image_fit(s, FIGS / "fig_sca_teaser_frontlight.png",
              Inches(7.6), Inches(1.7), Inches(5.3), Inches(5.0))
add_text(s, Inches(7.6), Inches(6.75), Inches(5.3), Inches(0.3),
         "Iso \u2192 A, Cow \u2192 S, Duck \u2192 C  \u00b7  single framework",
         size=10, color=MUTE, align=PP_ALIGN.CENTER)


# =====================================================================
# Slide 2 — Contributions
# =====================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Contributions", "Inject visual guidance where physics can absorb it")

# Three columns
col_w = Inches(4.0)
col_y = Inches(1.75)
col_h = Inches(5.0)
gap = Inches(0.15)
x0 = Inches(0.6)

cards = [
    ("1",
     "Control-space render guidance via MPM\u2013Gaussian duality",
     "Each MPM particle \u2261 one 3D Gaussian (\u03a3 = F \u03a3\u2080 F\u1d40). Route render gradients exclusively through \u2202L/\u2202F; set \u2202L/\u2202x = 0. Dimensionally aligned with Piola\u2013Kirchhoff stress \u2014 no cross-space scaling."),
    ("2",
     "Phased Chamfer-guided plasticity",
     "After a warmup k\u2080, update F_p via Chamfer NN displacements so the elastic rest state migrates toward the target. Restoring forces flip from adversary to ally; deformation persists across frames."),
    ("3",
     "Empirical analysis of render\u2013physics coupling",
     "At matched physics convergence, silhouette error drops 25.8% / 10.8% / 49.9% on Bunny / Cow / Duck. Largest gains on thin structures. Source-invariant: same target \u2192 shared attractor."),
]

for i, (num, title, body) in enumerate(cards):
    x = x0 + (col_w + gap) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_y, col_w, col_h)
    card.line.color.rgb = RULE
    card.line.width = Pt(0.75)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.shadow.inherit = False
    # Number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), col_y + Inches(0.3),
                               Inches(0.55), Inches(0.55))
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = "Calibri"; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_text(s, x + Inches(0.3), col_y + Inches(1.0), col_w - Inches(0.6), Inches(1.1),
             title, size=16, bold=True, color=INK)
    add_text(s, x + Inches(0.3), col_y + Inches(2.15), col_w - Inches(0.6), Inches(2.6),
             body, size=13, color=MUTE)


# =====================================================================
# Slide 3 — Method: MPM\u2013Gaussian Duality
# =====================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Method  \u00b7  1 of 4",
           "MPM\u2013Gaussian duality: one particle, two roles")

add_bullets(s, Inches(0.6), Inches(1.85), Inches(5.2), Inches(4.8), [
    ("Each MPM particle  \u2261  one 3D Gaussian primitive.", True),
    "Shared position:   \u03bc\u1d62 = x\u1d62",
    "Covariance from F:  \u03a3\u1d62 = F\u1d62 \u03a3\u2080 F\u1d62\u1d40,   \u03a3\u2080 = \u03c3\u2080\u00b2 I",
    ("F carries both stretch and rotation \u2192 rotation-gradient channel.", True),
    "When F = I, only stretch gradients exist; F \u2260 I unlocks rotation signal.",
    ("Fully differentiable \u2202L/\u2202F\u1d62 = (\u2202L/\u2202\u03a3\u1d62)(\u2202\u03a3\u1d62/\u2202F\u1d62).", True),
    "No explicit polar decomposition at render time (follows PhysGaussian).",
])

add_image_fit(s, FIGS / "fig_duality_v3.png",
              Inches(6.0), Inches(1.8), Inches(7.0), Inches(5.0))


# =====================================================================
# Slide 4 \u2014 Method: F-space Injection
# =====================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Method  \u00b7  2 of 4",
           "Control-space render injection: on F, not on x")

add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
         "(g^F,  g^x)  =  (\u03b3 \u00b7 \u2202L_render / \u2202F,   0)     with  \u03b3 = 0.1",
         size=20, bold=True, color=ACCENT2)

# Two columns: \u2202L/\u2202x (bad)  vs  \u2202L/\u2202F (good)
bad_x = Inches(0.6); good_x = Inches(7.0)
top = Inches(2.45); colw = Inches(5.7); colh = Inches(4.3)

for (x, title, color, items) in [
    (bad_x, "\u2202L/\u2202x  injection  \u2014  destabilizes", ACCENT, [
        "Lives in configuration space; MPM update is driven by stresses (F-space).",
        "No natural magnitude correspondence \u2192 requires ad-hoc cross-space scaling.",
        "Directly conflicts with elastic restoring forces on the same state variable.",
        "Empirically: eigen-decomposition failures at gains F absorbs without issue.",
        "Sparse coverage: nonzero only on surface-adjacent particles.",
    ]),
    (good_x, "\u2202L/\u2202F  injection  \u2014  composes with physics", ACCENT2, [
        "Dimensionally aligned with Piola\u2013Kirchhoff stress that P2G already integrates.",
        "Adjusts the control deformation field \u03b4F_c; physics produces the trajectory.",
        "Respects material properties, conservation, stability constraints.",
        "Broad coverage: propagates through every particle contributing to render pixels.",
        "Stabilized by drag d = 0.9 and temporal F-smoothing s = 0.955.",
    ]),
]:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, colw, colh)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.shadow.inherit = False
    add_text(s, x + Inches(0.3), top + Inches(0.25), colw - Inches(0.6), Inches(0.5),
             title, size=16, bold=True, color=color)
    add_bullets(s, x + Inches(0.3), top + Inches(0.85), colw - Inches(0.6), colh - Inches(1.0),
                items, size=13, color=INK)


# =====================================================================
# Slide 5 \u2014 Method: Chamfer-Guided Plasticity
# =====================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Method  \u00b7  3 of 4",
           "Chamfer-guided plasticity: migrate the elastic rest state")

add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
         "Problem: elastic energy \u03a8(F_e) has its minimum at F_e = I \u2192 achieved deformation rebounds toward source.",
         size=14, color=MUTE)

add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
         "Fix: multiplicative split  F = F_e F_p.  Update F_p so current configuration becomes the new rest state.",
         size=15, bold=True, color=INK)

add_bullets(s, Inches(0.6), Inches(2.95), Inches(12.2), Inches(4.0), [
    ("Activation:  only after warmup k\u2080 = 20 frames \u2014 let coarse structural seeds form first.", True),
    "(1)  NN displacement:   d\u1d62 = x\u1d62^nn \u2212 x\u1d62   (target surface, KD-tree).",
    "(2)  Smooth via iterative KNN averaging (k = 64, 3 iters) \u2014 suppresses NN noise.",
    "(3)  Symmetric Jacobian:   \u03b4F_p = \u00bd (J + J\u1d40),   J from KNN displacement gradient.",
    "(4)  Multiplicative update:   F_p \u2190 (I + \u03b7 \u03b4F_p) F_p,   \u03b7 = 0.2.",
    "(5)  Damp toward identity (d_p = 0.05) + isochoric projection \u2192 bounded det(F_p) & anisotropy.",
    ("Effect: restoring forces now point toward target \u2014 adversary becomes ally; rebound < 8%.", True),
])


# =====================================================================
# Slide 6 \u2014 Method: Pipeline overview
# =====================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Method  \u00b7  4 of 4",
           "Pipeline: three stages per frame, promoted across frames")

# Pipeline figure on top
add_image_fit(s, FIGS / "fig_pipeline_v4.png",
              Inches(0.6), Inches(1.55), Inches(12.1), Inches(4.45))

# Three labelled boxes / bullet summary beneath
stage_y = Inches(6.12)
for i, (title, body, col) in enumerate([
    ("A \u00b7 Physics rollout",
     "MPM forward T steps; Adam on \u03b4F_c; inject \u03b3 \u00b7 \u2202L/\u2202F from previous frame.",
     ACCENT2),
    ("B \u00b7 Multi-view render",
     "V = 8 cameras, 960\u00d7540; DT + soft IoU + depth; extract \u2202L_render/\u2202F.",
     ACCENT),
    ("C \u00b7 Chamfer plasticity",
     "After k \u2265 k\u2080: NN disp. \u2192 KNN smooth \u2192 F_p \u2190 (I + \u03b7 \u03b4F_p) F_p, damp + project.",
     RGBColor(0x2E, 0x8B, 0x57)),
]):
    x = Inches(0.6 + 4.2 * i)
    add_text(s, x, stage_y, Inches(4.0), Inches(0.35),
             title, size=13, bold=True, color=col)
    add_text(s, x, stage_y + Inches(0.35), Inches(4.0), Inches(0.9),
             body, size=11, color=INK)


prs.save(OUT)
print(f"wrote: {OUT}  ({OUT.stat().st_size/1024:.0f} KB)")
